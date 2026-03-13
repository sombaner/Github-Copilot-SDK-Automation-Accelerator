"""
PPT Generator Agent — Python

An AI-powered PowerPoint presentation generator built with the GitHub Copilot SDK.
Uses Claude Opus 4.6 to understand prompts, leverage agent skills, and invoke
custom tools to produce polished .pptx files.

Features:
- Accepts a user prompt describing the desired presentation
- Accepts agent skills as input to shape content & style
- Custom tools for slide creation, theming, chart insertion, and file export
- Uses Claude Opus 4.6 for intelligent content generation
"""

import asyncio
import json
import os
import re
import sys
from typing import Optional

from copilot import CopilotClient
from copilot.tools import define_tool
from copilot.generated.session_events import SessionEventType
from pydantic import BaseModel, Field
from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.chart import XL_CHART_TYPE

# ─── Types ────────────────────────────────────────────────────────────────────


class SlideContent(BaseModel):
    title: str
    body: Optional[str] = None
    bullets: Optional[list[str]] = None
    notes: Optional[str] = None
    layout: str = "content"  # title|content|section|two-column|comparison|chart|image
    chart_data: Optional[list[dict]] = None  # [{label: str, value: float}]
    left_column: Optional[str] = None
    right_column: Optional[str] = None
    image_url: Optional[str] = None
    image_caption: Optional[str] = None


class ThemeConfig(BaseModel):
    primary_color: str = "003366"
    secondary_color: str = "0066CC"
    accent_color: str = "FF6600"
    font_family: str = "Calibri"
    title_font_size: int = 28
    body_font_size: int = 16
    background_color: str = "FFFFFF"


# ─── In-Memory State ─────────────────────────────────────────────────────────

presentation_title = "Untitled Presentation"
slides: list[SlideContent] = []
current_theme = ThemeConfig()
agent_skills: list[str] = []


def _hex_to_rgb(hex_str: str) -> RGBColor:
    """Convert a hex colour string (without #) to an RGBColor."""
    hex_str = hex_str.lstrip("#")
    return RGBColor(int(hex_str[0:2], 16), int(hex_str[2:4], 16), int(hex_str[4:6], 16))


# ─── Custom Tools ─────────────────────────────────────────────────────────────


class SetPresentationMetadataParams(BaseModel):
    title: str = Field(description="Title of the presentation")
    skills: Optional[list[str]] = Field(
        default=None,
        description="Agent skills to apply (e.g. 'concise', 'technical', 'storytelling', 'data-driven', 'executive-summary')",
    )


@define_tool(description="Set the presentation title and optional agent skills. Call this first before adding slides.")
async def set_presentation_metadata(params: SetPresentationMetadataParams) -> dict:
    global presentation_title, agent_skills, slides
    presentation_title = params.title
    if params.skills:
        agent_skills = params.skills
    slides = []  # reset on new presentation
    return {
        "status": "metadata_set",
        "title": presentation_title,
        "skills": agent_skills,
        "message": f'Presentation "{presentation_title}" initialised with skills: [{", ".join(agent_skills)}]',
    }


class SetThemeParams(BaseModel):
    primary_color: Optional[str] = Field(default=None, description="Primary colour hex (without #)")
    secondary_color: Optional[str] = Field(default=None, description="Secondary colour hex")
    accent_color: Optional[str] = Field(default=None, description="Accent colour hex")
    font_family: Optional[str] = Field(default=None, description="Font family name")
    title_font_size: Optional[int] = Field(default=None, description="Title font size in points")
    body_font_size: Optional[int] = Field(default=None, description="Body font size in points")
    background_color: Optional[str] = Field(default=None, description="Background colour hex")


@define_tool(description="Set the visual theme for the presentation including colors, fonts, and sizes")
async def set_theme(params: SetThemeParams) -> dict:
    global current_theme
    updates = params.model_dump(exclude_none=True)
    current_theme = ThemeConfig(**{**current_theme.model_dump(), **updates})
    return {"status": "theme_updated", "theme": current_theme.model_dump()}


class AddSlideParams(BaseModel):
    title: str = Field(description="Slide title")
    body: Optional[str] = Field(default=None, description="Body text (content layout)")
    bullets: Optional[list[str]] = Field(default=None, description="Bullet points")
    notes: Optional[str] = Field(default=None, description="Speaker notes")
    layout: Optional[str] = Field(
        default="content",
        description="Layout: title, content, section, two-column, comparison, chart, image",
    )
    chart_data: Optional[list[dict]] = Field(
        default=None, description="Chart data [{label, value}]"
    )
    left_column: Optional[str] = Field(default=None, description="Left column text")
    right_column: Optional[str] = Field(default=None, description="Right column text")
    image_url: Optional[str] = Field(default=None, description="Image URL or path")
    image_caption: Optional[str] = Field(default=None, description="Image caption")


@define_tool(description="Add a slide. Layouts: title, content, section, two-column, comparison, chart, image.")
async def add_slide(params: AddSlideParams) -> dict:
    slide = SlideContent(**params.model_dump())
    slides.append(slide)
    return {
        "status": "slide_added",
        "slide_number": len(slides),
        "title": slide.title,
        "layout": slide.layout,
        "total_slides": len(slides),
    }


class EmptyParams(BaseModel):
    pass


@define_tool(description="List all slides currently in the presentation for review")
async def list_slides(params: EmptyParams) -> dict:
    return {
        "presentation_title": presentation_title,
        "theme": current_theme.model_dump(),
        "skills": agent_skills,
        "total_slides": len(slides),
        "slides": [
            {
                "number": i + 1,
                "title": s.title,
                "layout": s.layout,
                "has_bullets": bool(s.bullets),
                "has_chart": bool(s.chart_data),
                "has_notes": bool(s.notes),
            }
            for i, s in enumerate(slides)
        ],
    }


class RemoveSlideParams(BaseModel):
    slide_number: int = Field(description="Slide number to remove (1-based)")


@define_tool(description="Remove a slide by its slide number (1-based)")
async def remove_slide(params: RemoveSlideParams) -> dict:
    idx = params.slide_number - 1
    if idx < 0 or idx >= len(slides):
        return {"status": "error", "message": f"Invalid slide number {params.slide_number}. Total: {len(slides)}"}
    removed = slides.pop(idx)
    return {"status": "slide_removed", "removed_title": removed.title, "total_slides": len(slides)}


class ExportPptxParams(BaseModel):
    filename: Optional[str] = Field(default=None, description="Output filename (without extension)")
    output_dir: Optional[str] = Field(default="./output", description="Output directory path")


@define_tool(description="Export the presentation to a .pptx file. Call after all slides are added.")
async def export_pptx(params: ExportPptxParams) -> dict:
    if not slides:
        return {"status": "error", "message": "No slides to export. Add slides first."}

    prs = Presentation()
    prs.slide_width = Inches(10)
    prs.slide_height = Inches(5.625)  # 16:9

    blank_layout = prs.slide_layouts[6]  # blank

    for sc in slides:
        slide = prs.slides.add_slide(blank_layout)

        # Speaker notes
        if sc.notes:
            slide.notes_slide.notes_text_frame.text = sc.notes

        if sc.layout == "title":
            # Full-colour title slide
            bg = slide.background
            fill = bg.fill
            fill.solid()
            fill.fore_color.rgb = _hex_to_rgb(current_theme.primary_color)

            txbox = slide.shapes.add_textbox(Inches(0.5), Inches(1.5), Inches(9), Inches(1.5))
            tf = txbox.text_frame
            tf.word_wrap = True
            p = tf.paragraphs[0]
            p.text = sc.title
            p.font.size = Pt(current_theme.title_font_size + 8)
            p.font.bold = True
            p.font.color.rgb = RGBColor(0xFF, 0xFF, 0xFF)
            p.font.name = current_theme.font_family
            p.alignment = PP_ALIGN.CENTER

            if sc.body:
                txbox2 = slide.shapes.add_textbox(Inches(1.5), Inches(3.2), Inches(7), Inches(1))
                tf2 = txbox2.text_frame
                tf2.word_wrap = True
                p2 = tf2.paragraphs[0]
                p2.text = sc.body
                p2.font.size = Pt(current_theme.body_font_size + 4)
                p2.font.color.rgb = RGBColor(0xCC, 0xDD, 0xEE)
                p2.font.name = current_theme.font_family
                p2.alignment = PP_ALIGN.CENTER

        elif sc.layout == "section":
            bg = slide.background
            fill = bg.fill
            fill.solid()
            fill.fore_color.rgb = _hex_to_rgb(current_theme.secondary_color)

            txbox = slide.shapes.add_textbox(Inches(0.5), Inches(2), Inches(9), Inches(1.5))
            tf = txbox.text_frame
            tf.word_wrap = True
            p = tf.paragraphs[0]
            p.text = sc.title
            p.font.size = Pt(current_theme.title_font_size + 4)
            p.font.bold = True
            p.font.color.rgb = RGBColor(0xFF, 0xFF, 0xFF)
            p.font.name = current_theme.font_family
            p.alignment = PP_ALIGN.CENTER

        elif sc.layout in ("two-column", "comparison"):
            # Title
            txbox = slide.shapes.add_textbox(Inches(0.5), Inches(0.8), Inches(9), Inches(0.6))
            tf = txbox.text_frame
            tf.word_wrap = True
            p = tf.paragraphs[0]
            p.text = sc.title
            p.font.size = Pt(current_theme.title_font_size)
            p.font.bold = True
            p.font.color.rgb = _hex_to_rgb(current_theme.primary_color)
            p.font.name = current_theme.font_family

            # Left column
            txl = slide.shapes.add_textbox(Inches(0.5), Inches(1.6), Inches(4.2), Inches(3.5))
            tfl = txl.text_frame
            tfl.word_wrap = True
            tfl.paragraphs[0].text = sc.left_column or ""
            tfl.paragraphs[0].font.size = Pt(current_theme.body_font_size)
            tfl.paragraphs[0].font.name = current_theme.font_family

            # Right column
            txr = slide.shapes.add_textbox(Inches(5.3), Inches(1.6), Inches(4.2), Inches(3.5))
            tfr = txr.text_frame
            tfr.word_wrap = True
            tfr.paragraphs[0].text = sc.right_column or ""
            tfr.paragraphs[0].font.size = Pt(current_theme.body_font_size)
            tfr.paragraphs[0].font.name = current_theme.font_family

        elif sc.layout == "chart":
            # Title
            txbox = slide.shapes.add_textbox(Inches(0.5), Inches(0.8), Inches(9), Inches(0.6))
            tf = txbox.text_frame
            tf.word_wrap = True
            p = tf.paragraphs[0]
            p.text = sc.title
            p.font.size = Pt(current_theme.title_font_size)
            p.font.bold = True
            p.font.color.rgb = _hex_to_rgb(current_theme.primary_color)
            p.font.name = current_theme.font_family

            if sc.chart_data:
                from pptx.chart.data import CategoryChartData

                chart_data = CategoryChartData()
                chart_data.categories = [d.get("label", "") for d in sc.chart_data]
                chart_data.add_series(sc.title, [d.get("value", 0) for d in sc.chart_data])

                slide.shapes.add_chart(
                    XL_CHART_TYPE.COLUMN_CLUSTERED,
                    Inches(0.5), Inches(1.6), Inches(9), Inches(3.5),
                    chart_data,
                )

        elif sc.layout == "image":
            # Title
            txbox = slide.shapes.add_textbox(Inches(0.5), Inches(0.8), Inches(9), Inches(0.6))
            tf = txbox.text_frame
            tf.word_wrap = True
            p = tf.paragraphs[0]
            p.text = sc.title
            p.font.size = Pt(current_theme.title_font_size)
            p.font.bold = True
            p.font.color.rgb = _hex_to_rgb(current_theme.primary_color)
            p.font.name = current_theme.font_family

            # Placeholder for image
            caption_text = sc.image_caption or sc.image_url or "[Image placeholder]"
            txbox2 = slide.shapes.add_textbox(Inches(1.5), Inches(2.5), Inches(7), Inches(1))
            tf2 = txbox2.text_frame
            tf2.word_wrap = True
            p2 = tf2.paragraphs[0]
            p2.text = caption_text
            p2.font.size = Pt(current_theme.body_font_size)
            p2.font.color.rgb = RGBColor(0x99, 0x99, 0x99)
            p2.font.name = current_theme.font_family
            p2.alignment = PP_ALIGN.CENTER

        else:
            # Default content slide
            # Title
            txbox = slide.shapes.add_textbox(Inches(0.5), Inches(0.8), Inches(9), Inches(0.6))
            tf = txbox.text_frame
            tf.word_wrap = True
            p = tf.paragraphs[0]
            p.text = sc.title
            p.font.size = Pt(current_theme.title_font_size)
            p.font.bold = True
            p.font.color.rgb = _hex_to_rgb(current_theme.primary_color)
            p.font.name = current_theme.font_family

            if sc.bullets:
                txbox2 = slide.shapes.add_textbox(Inches(0.5), Inches(1.6), Inches(9), Inches(3.5))
                tf2 = txbox2.text_frame
                tf2.word_wrap = True
                for i, bullet in enumerate(sc.bullets):
                    if i == 0:
                        p2 = tf2.paragraphs[0]
                    else:
                        p2 = tf2.add_paragraph()
                    p2.text = f"• {bullet}"
                    p2.font.size = Pt(current_theme.body_font_size)
                    p2.font.name = current_theme.font_family
                    p2.font.color.rgb = RGBColor(0x33, 0x33, 0x33)
                    p2.space_before = Pt(8)
            elif sc.body:
                txbox2 = slide.shapes.add_textbox(Inches(0.5), Inches(1.6), Inches(9), Inches(3.5))
                tf2 = txbox2.text_frame
                tf2.word_wrap = True
                p2 = tf2.paragraphs[0]
                p2.text = sc.body
                p2.font.size = Pt(current_theme.body_font_size)
                p2.font.name = current_theme.font_family
                p2.font.color.rgb = RGBColor(0x33, 0x33, 0x33)

    # Write file
    out_dir = params.output_dir or "./output"
    os.makedirs(out_dir, exist_ok=True)
    safe_name = re.sub(r"[^a-zA-Z0-9_-]", "_", params.filename or presentation_title)[:100]
    file_path = os.path.join(out_dir, f"{safe_name}.pptx")
    prs.save(file_path)

    return {
        "status": "exported",
        "file_path": file_path,
        "total_slides": len(slides),
        "title": presentation_title,
        "message": f"Presentation exported to {file_path} ({len(slides)} slides)",
    }


# ─── Main ─────────────────────────────────────────────────────────────────────


async def main():
    print("📊 PPT Generator Agent starting...\n")

    client = CopilotClient()
    await client.start()

    # Parse --skills from CLI
    input_skills: list[str] = []
    for arg in sys.argv[1:]:
        if arg.startswith("--skills="):
            input_skills = [s.strip() for s in arg.split("=", 1)[1].split(",")]

    skills_prompt = ""
    if input_skills:
        skills_prompt = f"""

Agent skills active for this session: [{", ".join(input_skills)}].
Adapt your content generation style according to these skills:
- "concise": keep text short and punchy, favour bullets over paragraphs
- "technical": use precise terminology, include code snippets or architecture diagrams where relevant
- "storytelling": structure the presentation as a narrative arc with intro, conflict, resolution
- "data-driven": emphasise numbers, charts, and evidence; use the chart slide layout when data is available
- "executive-summary": high-level overview, limit to 6-8 slides, focus on outcomes and recommendations
- "creative": use bold colour themes, metaphors, and visual variety
- Any other skill: interpret it as a content-styling directive and apply accordingly."""

    session = await client.create_session({
        "model": "claude-opus-4.6",
        "streaming": True,
        "tools": [
            set_presentation_metadata,
            set_theme,
            add_slide,
            list_slides,
            remove_slide,
            export_pptx,
        ],
        "system_message": {
            "content": f"""You are an expert presentation designer and content strategist. Your role is to:

1. Understand the user's prompt and create a professional PowerPoint presentation.
2. Use the provided tools to build the presentation step by step:
   a. First call set_presentation_metadata to set the title and any active skills.
   b. Optionally call set_theme to customise colours/fonts if the user requests a specific look.
   c. Call add_slide repeatedly to build out each slide with the appropriate layout.
   d. Call list_slides to review the deck before exporting.
   e. Call export_pptx to generate the final .pptx file.
3. Choose the best slide layout for each piece of content:
   - "title" for the opening slide
   - "section" for section dividers
   - "content" for standard bullet or text slides
   - "two-column" or "comparison" for side-by-side information
   - "chart" when presenting numerical data
   - "image" when referencing visuals
4. Always add speaker notes to every slide.
5. Apply the active agent skills to shape content tone, depth, and structure.
6. Keep individual slides focused — one key idea per slide.
7. End with a summary or call-to-action slide.
{skills_prompt}""",
        },
    })

    def handle_event(event):
        if event.type == SessionEventType.ASSISTANT_MESSAGE_DELTA:
            sys.stdout.write(event.data.delta_content)
            sys.stdout.flush()
        elif event.type == SessionEventType.SESSION_IDLE:
            print()

    session.on(handle_event)

    print("🚀 PPT Generator Agent Ready!")
    print("=" * 60)
    print(f"   Model  : claude-opus-4.6")
    if input_skills:
        print(f"   Skills : {', '.join(input_skills)}")
    print(f"   Output : ./output/")
    print("=" * 60)
    print("\nDescribe the presentation you want to create.")
    print("Type 'exit' to quit.\n")

    while True:
        try:
            user_input = input("You: ")
        except EOFError:
            break

        if user_input.lower() == "exit":
            break

        if not user_input.strip():
            continue

        sys.stdout.write("\nAssistant: ")
        await session.send_and_wait({"prompt": user_input})
        print("\n")

    print("\n👋 Shutting down PPT Generator Agent...")
    await client.stop()


if __name__ == "__main__":
    asyncio.run(main())
